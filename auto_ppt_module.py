# -*- coding: utf-8 -*-
"""
@author: ChaoTzuYin
"""
from pptx import Presentation
import numpy as np
import io
from PIL import Image
import copy



class info_keeper():
    def __init__(self, slide_idx, shape_idx, left, top, width, height):
        self.slide_idx = slide_idx  
        self.shape_idx = shape_idx
        self.left = left 
        self.top = top 
        self.width = width 
        self.height = height
        

class ppt_recorder():
    
    def __init__(self, template):
        
        '''
        template [string]: path to your template pptx file
        '''
        
        self.prs_ori = Presentation(template)
        self.prs = Presentation(template)

        self.bias = -len(self.prs_ori.slides)
        
        self.variables = {}
        for count in range(len(self.prs_ori.slides)):
            slide = self.prs_ori.slides[count]
            for shape_count in range(len(slide.shapes)):
                shape = slide.shapes[shape_count]
                
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    if('#(' in text_frame.text):
                       variable_name = text_frame.text[2:-1]
                       self.variables[variable_name] = info_keeper(slide_idx=count,
                                                                    shape_idx=shape_count,
                                                                    left=shape.left, 
                                                                    top=shape.top, 
                                                                    width=shape.width, 
                                                                    height=shape.height)

    def _get_blank_slide_layout(self,pres):
        layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return pres.slide_layouts[blank_layout_id]
    
    def new_record(self):
        """Duplicate the slide with the given index in pres.
    
        Adds slide to the end of the presentation"""
        if self.bias != -len(self.prs_ori.slides):
            for source in self.prs_ori.slides:
                blank_slide_layout = self._get_blank_slide_layout(self.prs)
                dest = self.prs.slides.add_slide(blank_slide_layout)
            
                for shp in source.shapes:
                    el = shp.element
                    newel = copy.deepcopy(el)
                    dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        
        self.bias = self.bias + len(self.prs_ori.slides)
    
    def placeholder(self):
        return self.variables

    
    def assign(self, feed_dict):
        
        assert self.bias != -len(self.prs_ori.slides), 'NULL SLIDE ERROR: Please call the function "new_record()" to create a new set of slides, with the formats respected to the template, before recording.'
        
        for holder, data in feed_dict.items():
            if(isinstance(data, np.ndarray)):
                if(len(data.shape)==2):
                    np_image = np.tile(data[:,:,None]*255,(1,1,3)).astype(np.uint8)
                else:
                    np_image = (data*255).astype(np.uint8)
                
                s = io.BytesIO()
                Image.fromarray(np_image).save(s, format='png')
                s.seek(0)
                x, y, cx, cy = holder.left, holder.top, holder.width, holder.height
                Image.MAX_IMAGE_PIXELS = 1000000000
                self.prs.slides[holder.slide_idx+self.bias].shapes.add_picture(io.BytesIO(s.getvalue()), x, y, cx, cy)
                s.truncate(0)
                s.close()
            
            elif(isinstance(data,str)):
                
                text_frame = self.prs.slides[holder.slide_idx+self.bias].shapes[holder.shape_idx].text_frame
                
                cur_text = text_frame.text
                new_text = cur_text.replace(cur_text, data)
                text_frame.text = new_text
            
            else:
                print(type(data))
    
    def to_pptx(self, path):
        self.prs.save(path)

    
    
if __name__ == '__main__':

    # Create ppt recorder given the reference template.
    writer = ppt_recorder(template='template_example.pptx')
    
    # Get the placeholders that you've created in the template file. 
    # Placeholders are formated in dictionary, in which the keys are repected to their name.
    
    ph = writer.placeholder()
    
    test_image1 = np.eye(128)
    test_image2 = np.random.uniform(low=0.0, high=1.0, size=[64,64,3])

    for i in range(3):
        # Duplicate a set of formate respected to your template pptx file.
        # You must include .new_record() in your code even if only one set of your template is needed in your report.
        writer.new_record() 
        
        # Create a placeholder-data dictionary.
        feed_dict={ph['text1']:'Hello World!_'+str(i), 
                 ph['pic1']:test_image1,
                 ph['pic2']:test_image2,
                 ph['text2']:"What's up, World?_"+str(i),
                 ph['pic3']:test_image1,
                 ph['pic4']:test_image2}
        
        # Assign your data to pptx file.
        writer.assign(feed_dict=feed_dict)
    
    #export pptx file.
    writer.to_pptx('result.pptx')

